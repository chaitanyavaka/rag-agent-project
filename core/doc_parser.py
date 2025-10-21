import os
import pypdf
import docx
import pptx
import pandas as pd

def parse_document(file_path):
    """Parses a document and extracts text based on its file extension."""
    _, extension = os.path.splitext(file_path)
    text = ""
    
    try:
        if extension.lower() == '.pdf':
            with open(file_path, 'rb') as f:
                reader = pypdf.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        
        elif extension.lower() == '.docx':
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        
        elif extension.lower() == '.pptx':
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        elif extension.lower() == '.csv':
            df = pd.read_csv(file_path)
            text += df.to_string()
        
        elif extension.lower() in ['.txt', '.md']:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
    
    except Exception as e:
        print(f"Error parsing {file_path}: {e}")
        return f"Error parsing file: {os.path.basename(file_path)}"
        
    return text